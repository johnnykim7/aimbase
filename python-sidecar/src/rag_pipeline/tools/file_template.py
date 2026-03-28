"""CR-018: File-based Document Template Engine.

PPTX/DOCX 파일을 템플릿으로 저장하고, {{변수}} 플레이스홀더를 치환하여 문서 생성.
- upload: base64 파일 → 로컬 저장 + {{변수}} 자동 추출 → DB 등록
- render: 파일 로드 → 변수 치환 → base64 결과 반환
"""

import base64
import json
import logging
import os
import re
import subprocess
import tempfile
import uuid
from typing import Any

logger = logging.getLogger(__name__)

# LibreOffice 경로 (Mac / Linux)
LIBREOFFICE_PATH = os.getenv("LIBREOFFICE_PATH", "")
_LO_CANDIDATES = [
    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    "/usr/bin/libreoffice",
    "/usr/bin/soffice",
]

# 템플릿 파일 저장 경로
TEMPLATE_STORAGE_DIR = os.getenv("TEMPLATE_STORAGE_DIR", "./data/templates")

PLACEHOLDER_PATTERN = re.compile(r"\{\{(\w+)\}\}")


def _find_libreoffice() -> str:
    """LibreOffice 실행 파일 경로 탐색."""
    if LIBREOFFICE_PATH and os.path.exists(LIBREOFFICE_PATH):
        return LIBREOFFICE_PATH
    for candidate in _LO_CANDIDATES:
        if os.path.exists(candidate):
            return candidate
    return ""


def _convert_with_libreoffice(input_path: str, output_format: str, output_dir: str) -> str:
    """LibreOffice headless로 포맷 변환. 변환된 파일 경로 반환."""
    lo_path = _find_libreoffice()
    if not lo_path:
        raise RuntimeError("LibreOffice not found. Install it for format conversion.")

    cmd = [lo_path, "--headless", "--convert-to", output_format, "--outdir", output_dir, input_path]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)

    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

    # 변환된 파일 찾기
    basename = os.path.splitext(os.path.basename(input_path))[0]
    converted_path = os.path.join(output_dir, f"{basename}.{output_format}")
    if os.path.exists(converted_path):
        return converted_path

    # 파일명이 다를 수 있으므로 디렉토리에서 탐색
    for f in os.listdir(output_dir):
        if f.endswith(f".{output_format}"):
            return os.path.join(output_dir, f)

    raise RuntimeError(f"Converted file not found in {output_dir}")


def _convert_pdf_to_docx(pdf_path: str, docx_path: str):
    """PDF → DOCX 변환 (pdf2docx 사용)."""
    from pdf2docx import Converter
    cv = Converter(pdf_path)
    cv.convert(docx_path)
    cv.close()


def _ensure_storage_dir():
    os.makedirs(TEMPLATE_STORAGE_DIR, exist_ok=True)


# ── PPTX 처리 ────────────────────────────────────────────────

def _extract_variables_pptx(file_path: str) -> list[str]:
    """PPTX 파일의 모든 텍스트에서 {{변수명}} 추출."""
    from pptx import Presentation
    prs = Presentation(file_path)
    found = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    full_text = "".join(run.text for run in para.runs)
                    found.update(PLACEHOLDER_PATTERN.findall(full_text))
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            full_text = "".join(run.text for run in para.runs)
                            found.update(PLACEHOLDER_PATTERN.findall(full_text))
    return sorted(found)


def _render_pptx(file_path: str, variables: dict[str, str], output_path: str):
    """PPTX 파일의 {{변수}}를 치환하여 새 파일 저장."""
    from pptx import Presentation
    prs = Presentation(file_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                _replace_in_text_frame(shape.text_frame, variables)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        _replace_in_text_frame(cell.text_frame, variables)

    prs.save(output_path)


def _replace_in_text_frame(text_frame, variables: dict[str, str]):
    """TextFrame 내 runs를 합쳐서 {{변수}} 치환 후 첫 번째 run에 결과를 넣고 나머지 비움.

    python-pptx는 {{title}}이 여러 run으로 쪼개질 수 있으므로
    paragraph 단위로 전체 텍스트를 합쳐서 치환합니다.
    """
    for para in text_frame.paragraphs:
        runs = para.runs
        if not runs:
            continue

        full_text = "".join(run.text for run in runs)
        if "{{" not in full_text:
            continue

        new_text = full_text
        for key, value in variables.items():
            new_text = new_text.replace("{{" + key + "}}", str(value))

        # 첫 번째 run에 결과를 넣고 나머지 run은 비움
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""


# ── DOCX 처리 ────────────────────────────────────────────────

def _extract_variables_docx(file_path: str) -> list[str]:
    """DOCX 파일의 모든 텍스트에서 {{변수명}} 추출."""
    from docx import Document
    doc = Document(file_path)
    found = set()

    for para in doc.paragraphs:
        full_text = "".join(run.text for run in para.runs)
        found.update(PLACEHOLDER_PATTERN.findall(full_text))

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    full_text = "".join(run.text for run in para.runs)
                    found.update(PLACEHOLDER_PATTERN.findall(full_text))

    # 헤더/푸터
    for section in doc.sections:
        for header_footer in [section.header, section.footer]:
            if header_footer is not None:
                for para in header_footer.paragraphs:
                    full_text = "".join(run.text for run in para.runs)
                    found.update(PLACEHOLDER_PATTERN.findall(full_text))

    return sorted(found)


def _render_docx(file_path: str, variables: dict[str, str], output_path: str):
    """DOCX 파일의 {{변수}}를 치환하여 새 파일 저장."""
    from docx import Document
    doc = Document(file_path)

    for para in doc.paragraphs:
        _replace_in_docx_paragraph(para, variables)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _replace_in_docx_paragraph(para, variables)

    for section in doc.sections:
        for header_footer in [section.header, section.footer]:
            if header_footer is not None:
                for para in header_footer.paragraphs:
                    _replace_in_docx_paragraph(para, variables)

    doc.save(output_path)


def _replace_in_docx_paragraph(para, variables: dict[str, str]):
    """DOCX paragraph 내 runs를 합쳐서 {{변수}} 치환."""
    runs = para.runs
    if not runs:
        return

    full_text = "".join(run.text for run in runs)
    if "{{" not in full_text:
        return

    new_text = full_text
    for key, value in variables.items():
        new_text = new_text.replace("{{" + key + "}}", str(value))

    runs[0].text = new_text
    for run in runs[1:]:
        run.text = ""


# ── 공개 API ─────────────────────────────────────────────────

def upload_file_template(
    name: str,
    format: str,
    file_base64: str,
    original_filename: str = "",
    description: str = "",
    tags: str = "[]",
) -> dict[str, Any]:
    """파일을 업로드하여 템플릿으로 등록.

    Args:
        name: 템플릿 이름
        format: 파일 포맷 (pptx, docx)
        file_base64: base64 인코딩된 파일 내용
        original_filename: 원본 파일명
        description: 설명
        tags: JSON 배열 문자열

    Returns:
        {template_id, name, variables, success}
    """
    supported = ("pptx", "docx", "pdf")
    if format not in supported:
        return {"success": False, "error": f"File template supports {supported}, got: {format}"}

    _ensure_storage_dir()

    template_id = str(uuid.uuid4())
    file_bytes = base64.b64decode(file_base64)

    # PDF인 경우 DOCX로 변환하여 내부 저장
    original_format = format
    if format == "pdf":
        try:
            pdf_tmp = os.path.join(TEMPLATE_STORAGE_DIR, f"{template_id}.pdf")
            with open(pdf_tmp, "wb") as f:
                f.write(file_bytes)
            stored_path = os.path.join(TEMPLATE_STORAGE_DIR, f"{template_id}.docx")
            _convert_pdf_to_docx(pdf_tmp, stored_path)
            os.remove(pdf_tmp)
            format = "docx"  # 내부적으로 docx로 관리
            logger.info("PDF converted to DOCX for template storage")
        except Exception as e:
            return {"success": False, "error": f"PDF to DOCX conversion failed: {e}"}
    else:
        ext = f".{format}"
        stored_filename = f"{template_id}{ext}"
        stored_path = os.path.join(TEMPLATE_STORAGE_DIR, stored_filename)
        with open(stored_path, "wb") as f:
            f.write(file_bytes)

    # 변수 추출
    try:
        if format == "pptx":
            var_names = _extract_variables_pptx(stored_path)
        else:
            var_names = _extract_variables_docx(stored_path)
    except Exception as e:
        os.remove(stored_path)
        return {"success": False, "error": f"Failed to parse file: {e}"}

    variables = [
        {"name": v, "type": "string", "default_value": "", "required": True, "description": ""}
        for v in var_names
    ]

    # tags 파싱
    try:
        tags_list = json.loads(tags) if isinstance(tags, str) else tags
    except json.JSONDecodeError:
        tags_list = []

    # DB 저장
    from rag_pipeline.db import get_connection
    conn = get_connection()
    try:
        with conn.cursor() as cur:
            cur.execute(
                """
                INSERT INTO document_templates
                    (id, name, description, format, template_type, file_path, variables, tags)
                VALUES (%s, %s, %s, %s, 'file', %s, %s::jsonb, %s::jsonb)
                """,
                (
                    template_id, name, description, format,
                    stored_path,
                    json.dumps(variables, ensure_ascii=False),
                    json.dumps(tags_list, ensure_ascii=False),
                ),
            )
        conn.commit()

        logger.info("File template uploaded: %s (%s→%s), %d variables", name, original_format, format, len(variables))
        return {
            "template_id": template_id,
            "name": name,
            "format": format,
            "original_format": original_format,
            "converted": original_format != format,
            "original_filename": original_filename,
            "variables": variables,
            "variable_count": len(variables),
            "file_size_bytes": len(file_bytes),
            "success": True,
        }
    except Exception as e:
        conn.rollback()
        os.remove(stored_path)
        return {"success": False, "error": str(e)}


def render_file_template(
    template_id: str,
    variables: str = "{}",
    output_format: str = "",
) -> dict[str, Any]:
    """파일 템플릿의 {{변수}}를 치환하여 문서 생성.

    Args:
        template_id: 템플릿 ID
        variables: JSON 객체 문자열
        output_format: 출력 포맷 (빈 문자열이면 템플릿 원본 포맷, pdf/docx/pptx 지정 가능)

    Returns:
        {file_base64, filename, format, size_bytes, success}
    """
    try:
        var_values = json.loads(variables) if isinstance(variables, str) else variables
    except json.JSONDecodeError:
        return {"success": False, "error": "Invalid variables JSON"}

    # 템플릿 로드
    from rag_pipeline.tools.template_manager import get_template
    tpl_result = get_template(template_id)
    if not tpl_result.get("success"):
        return tpl_result

    tpl = tpl_result["template"]

    if tpl["template_type"] != "file":
        return {"success": False, "error": "Not a file template. Use render_document_template for code templates."}

    file_path = tpl.get("file_path", "")
    if not file_path or not os.path.exists(file_path):
        return {"success": False, "error": f"Template file not found: {file_path}"}

    internal_fmt = tpl["format"]  # 내부 저장 포맷 (docx or pptx)
    target_fmt = output_format.lower().strip() if output_format else internal_fmt

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            # 1단계: 변수 치환 (내부 포맷으로)
            rendered_filename = f"{tpl['name']}.{internal_fmt}"
            rendered_path = os.path.join(tmpdir, rendered_filename)

            if internal_fmt == "pptx":
                _render_pptx(file_path, var_values, rendered_path)
            elif internal_fmt == "docx":
                _render_docx(file_path, var_values, rendered_path)
            else:
                return {"success": False, "error": f"File rendering not supported for format: {internal_fmt}"}

            # 2단계: 출력 포맷 변환 (필요한 경우)
            if target_fmt != internal_fmt:
                try:
                    final_path = _convert_with_libreoffice(rendered_path, target_fmt, tmpdir)
                    output_filename = f"{tpl['name']}.{target_fmt}"
                except RuntimeError as e:
                    return {"success": False, "error": str(e)}
            else:
                final_path = rendered_path
                output_filename = rendered_filename

            with open(final_path, "rb") as f:
                content = f.read()

            file_base64 = base64.b64encode(content).decode("utf-8")
            size_bytes = len(content)

            logger.info("File template rendered: %s (%d bytes)", output_filename, size_bytes)
            return {
                "file_base64": file_base64,
                "filename": output_filename,
                "format": target_fmt,
                "size_bytes": size_bytes,
                "template_id": template_id,
                "template_name": tpl["name"],
                "success": True,
            }
    except Exception as e:
        logger.error("File template rendering failed: %s", e)
        return {"success": False, "error": f"Rendering failed: {e}"}
