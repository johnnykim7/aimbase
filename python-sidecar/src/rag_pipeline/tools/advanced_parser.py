"""PY-027: Advanced Document Parsing MCP Tool (CR-011).

기존 parser.py(PY-013)의 확장판.
- Unstructured 라이브러리 기반 고급 파싱 (테이블/이미지/레이아웃 인식)
- 문서 구조(헤딩, 섹션) 보존
- 메타데이터 강화 (페이지 번호, 섹션 제목, 요소 타입)
- Fallback: PyMuPDF / regex 기반 파싱
"""

import base64
import io
import logging
import re
from typing import Any

logger = logging.getLogger(__name__)

# Unstructured 사용 가능 여부 확인
_HAS_UNSTRUCTURED = False
try:
    from unstructured.partition.auto import partition
    from unstructured.partition.pdf import partition_pdf
    from unstructured.partition.docx import partition_docx
    from unstructured.partition.pptx import partition_pptx
    from unstructured.partition.html import partition_html
    from unstructured.partition.md import partition_md
    from unstructured.partition.text import partition_text
    _HAS_UNSTRUCTURED = True
    logger.info("Unstructured library loaded — advanced parsing enabled")
except ImportError:
    logger.warning("Unstructured not installed, falling back to basic parsing")


def advanced_parse(
    file_content: str,
    file_type: str = "",
    config: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """고급 문서 파싱 — 구조 보존 + 메타데이터 추출.

    Unstructured 라이브러리가 설치되어 있으면 자동으로 사용.
    없으면 PyMuPDF / regex 기반 fallback.

    Args:
        file_content: Base64 인코딩된 파일 바이트 또는 텍스트
        file_type: 파일 타입 힌트 ("pdf", "docx", "pptx", "html", "md", "txt")
        config: {extract_tables, extract_images, preserve_structure, use_unstructured}

    Returns:
        {sections: [{title, content, page?, metadata}], full_text, tables, success, parser}
    """
    cfg = config or {}
    extract_tables = cfg.get("extract_tables", True)
    preserve_structure = cfg.get("preserve_structure", True)
    use_unstructured = cfg.get("use_unstructured", True)

    try:
        raw_bytes = base64.b64decode(file_content)
    except Exception:
        raw_bytes = None
        raw_text = file_content

    ft = file_type.lower().strip() if file_type else ""

    # Unstructured 우선 시도
    if _HAS_UNSTRUCTURED and use_unstructured and ft in ("pdf", "docx", "pptx", "html", "xlsx"):
        try:
            return _parse_with_unstructured(raw_bytes, raw_text if raw_bytes is None else None, ft, extract_tables, preserve_structure)
        except Exception as e:
            logger.warning("Unstructured parsing failed for '%s', falling back: %s", ft, e)

    # Fallback: 기존 파서
    if ft in ("md", "markdown") or (raw_bytes is None and not ft):
        text = raw_text if raw_bytes is None else raw_bytes.decode("utf-8", errors="replace")
        return _parse_markdown(text, preserve_structure)

    if ft == "html":
        text = raw_bytes.decode("utf-8", errors="replace") if raw_bytes else raw_text
        return _parse_html(text, preserve_structure, extract_tables)

    if ft == "txt":
        text = raw_bytes.decode("utf-8", errors="replace") if raw_bytes else raw_text
        return _parse_plain_text(text, preserve_structure)

    if ft == "pdf" and raw_bytes:
        return _parse_pdf(raw_bytes, preserve_structure, extract_tables)

    if ft in ("docx", "pptx") and raw_bytes:
        return _parse_office_fallback(raw_bytes, ft, preserve_structure, extract_tables)

    # fallback: 텍스트로 처리
    text = raw_bytes.decode("utf-8", errors="replace") if raw_bytes else raw_text
    return _parse_plain_text(text, preserve_structure)


def _parse_with_unstructured(
    raw_bytes: bytes | None,
    raw_text: str | None,
    file_type: str,
    extract_tables: bool,
    preserve_structure: bool,
) -> dict[str, Any]:
    """Unstructured 라이브러리를 사용한 고급 파싱.

    테이블/이미지 캡션/레이아웃을 자동 인식합니다.
    """
    import tempfile
    import os

    # 파일 타입별 파티셔너 매핑
    partitioner_map = {
        "pdf": partition_pdf,
        "docx": partition_docx,
        "pptx": partition_pptx,
        "html": partition_html,
    }

    partitioner = partitioner_map.get(file_type, partition)

    # Unstructured는 파일 경로 또는 파일 객체가 필요
    if raw_bytes:
        suffix = f".{file_type}" if file_type else ""
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
            tmp.write(raw_bytes)
            tmp_path = tmp.name
        try:
            elements = partitioner(filename=tmp_path)
        finally:
            os.unlink(tmp_path)
    elif raw_text:
        elements = partition_text(text=raw_text)
    else:
        return {"sections": [], "full_text": "", "tables": [], "success": False, "error": "No content"}

    # 요소를 섹션/테이블로 분류
    sections: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []
    full_text_parts: list[str] = []
    current_section_title = ""
    current_section_content: list[str] = []

    for el in elements:
        el_type = type(el).__name__
        el_text = str(el).strip()
        el_meta = el.metadata if hasattr(el, "metadata") else None
        page_num = getattr(el_meta, "page_number", None) if el_meta else None

        if not el_text:
            continue

        full_text_parts.append(el_text)

        if el_type == "Title":
            # 이전 섹션 저장
            if current_section_title or current_section_content:
                sections.append({
                    "title": current_section_title,
                    "content": "\n".join(current_section_content).strip(),
                    "level": 0,
                    "metadata": {"parser": "unstructured"},
                })
            current_section_title = el_text
            current_section_content = []

        elif el_type == "Table" and extract_tables:
            # 테이블 요소
            table_data = {"content": el_text, "page": page_num}
            if hasattr(el_meta, "text_as_html") and el_meta.text_as_html:
                table_data["html"] = el_meta.text_as_html
                # HTML 테이블에서 행/셀 추출
                rows = _extract_table_rows_from_html(el_meta.text_as_html)
                if rows:
                    table_data["rows"] = rows
            tables.append(table_data)
            current_section_content.append(f"[표]\n{el_text}")

        elif el_type == "Image":
            caption = el_text if el_text else "[이미지]"
            current_section_content.append(f"[이미지: {caption}]")

        elif el_type == "FigureCaption":
            current_section_content.append(f"[캡션: {el_text}]")

        elif el_type == "Header":
            current_section_content.append(el_text)

        elif el_type == "Footer":
            pass  # 푸터 무시

        elif el_type == "PageBreak":
            pass  # 페이지 구분 무시

        else:
            # NarrativeText, ListItem, Address 등
            if page_num and preserve_structure:
                current_section_content.append(el_text)
            else:
                current_section_content.append(el_text)

    # 마지막 섹션 저장
    if current_section_title or current_section_content:
        sections.append({
            "title": current_section_title,
            "content": "\n".join(current_section_content).strip(),
            "level": 0,
            "metadata": {"parser": "unstructured"},
        })

    full_text = "\n\n".join(full_text_parts)

    logger.info("Unstructured parsed %s: %d elements → %d sections, %d tables",
                file_type, len(elements), len(sections), len(tables))

    return {
        "sections": sections,
        "full_text": full_text,
        "tables": tables,
        "format": file_type,
        "section_count": len(sections),
        "table_count": len(tables),
        "element_count": len(elements),
        "parser": "unstructured",
        "success": True,
    }


def _extract_table_rows_from_html(html: str) -> list[list[str]]:
    """Unstructured의 text_as_html에서 행/셀 데이터를 추출."""
    rows = re.findall(r"<tr[^>]*>(.*?)</tr>", html, re.DOTALL | re.IGNORECASE)
    result = []
    for row in rows:
        cells = re.findall(r"<t[dh][^>]*>(.*?)</t[dh]>", row, re.DOTALL | re.IGNORECASE)
        result.append([re.sub(r"<[^>]+>", "", c).strip() for c in cells])
    return result


def _parse_office_fallback(
    raw_bytes: bytes, file_type: str, preserve_structure: bool, extract_tables: bool,
) -> dict[str, Any]:
    """DOCX/PPTX — python-docx/python-pptx 기반 fallback."""
    if file_type == "docx":
        try:
            from docx import Document
            doc = Document(io.BytesIO(raw_bytes))
            sections = []
            full_text_parts = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                full_text_parts.append(text)
                if para.style and para.style.name.startswith("Heading"):
                    sections.append({"title": text, "content": "", "level": 0})
                elif sections:
                    sections[-1]["content"] += text + "\n"
                else:
                    sections.append({"title": "", "content": text + "\n", "level": 0})
            return {
                "sections": sections,
                "full_text": "\n".join(full_text_parts),
                "tables": [],
                "format": "docx",
                "parser": "python-docx",
                "section_count": len(sections),
                "success": True,
            }
        except ImportError:
            pass

    if file_type == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw_bytes))
            sections = []
            tables = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_parts = []
                for shape in slide.shapes:
                    if shape.has_table and extract_tables:
                        table = shape.table
                        rows = []
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            rows.append(cells)
                        tables.append({"page": slide_idx, "rows": rows})
                        slide_parts.append("[표]\n" + "\n".join(" | ".join(r) for r in rows))
                    elif hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text)
                if slide_parts:
                    sections.append({
                        "title": f"슬라이드 {slide_idx}",
                        "content": "\n".join(slide_parts),
                        "page": slide_idx,
                        "level": 0,
                    })
            full_text = "\n\n".join(s["content"] for s in sections)
            return {
                "sections": sections,
                "full_text": full_text,
                "tables": tables,
                "format": "pptx",
                "parser": "python-pptx",
                "section_count": len(sections),
                "table_count": len(tables),
                "success": True,
            }
        except ImportError:
            pass

    return {"sections": [], "full_text": "", "tables": [], "success": False, "error": f"No parser for {file_type}"}


def _parse_markdown(text: str, preserve_structure: bool) -> dict[str, Any]:
    """마크다운 구조 파싱 — 헤딩 기반 섹션 분할."""
    sections = []
    current_title = ""
    current_content: list[str] = []
    tables: list[dict] = []

    for line in text.split("\n"):
        heading_match = re.match(r"^(#{1,6})\s+(.+)", line)
        if heading_match and preserve_structure:
            # 이전 섹션 저장
            if current_title or current_content:
                sections.append({
                    "title": current_title,
                    "content": "\n".join(current_content).strip(),
                    "level": 0,
                })
            current_title = heading_match.group(2).strip()
            current_content = []
        elif line.strip().startswith("|") and "|" in line[1:]:
            # 테이블 행 감지
            if not tables or not current_content or current_content[-1].startswith("|"):
                current_content.append(line)
            else:
                current_content.append(line)
        else:
            current_content.append(line)

    # 마지막 섹션
    if current_title or current_content:
        sections.append({
            "title": current_title,
            "content": "\n".join(current_content).strip(),
            "level": 0,
        })

    return {
        "sections": sections,
        "full_text": text,
        "tables": tables,
        "format": "markdown",
        "parser": "regex",
        "section_count": len(sections),
        "success": True,
    }


def _parse_html(text: str, preserve_structure: bool, extract_tables: bool) -> dict[str, Any]:
    """HTML 구조 파싱."""
    # 간단한 태그 제거 + 구조 보존
    sections = []
    tables = []

    # 테이블 추출
    if extract_tables:
        table_pattern = re.compile(r"<table[^>]*>(.*?)</table>", re.DOTALL | re.IGNORECASE)
        for match in table_pattern.finditer(text):
            rows = re.findall(r"<tr[^>]*>(.*?)</tr>", match.group(1), re.DOTALL | re.IGNORECASE)
            table_rows = []
            for row in rows:
                cells = re.findall(r"<t[dh][^>]*>(.*?)</t[dh]>", row, re.DOTALL | re.IGNORECASE)
                table_rows.append([re.sub(r"<[^>]+>", "", c).strip() for c in cells])
            if table_rows:
                tables.append({"rows": table_rows})

    # 헤딩 기반 섹션 분할
    heading_pattern = re.compile(r"<h([1-6])[^>]*>(.*?)</h\1>", re.DOTALL | re.IGNORECASE)
    parts = heading_pattern.split(text)

    # 태그 정리
    clean_text = re.sub(r"<[^>]+>", " ", text)
    clean_text = re.sub(r"\s+", " ", clean_text).strip()

    if preserve_structure and len(parts) > 1:
        for i in range(1, len(parts), 3):
            if i + 1 < len(parts):
                level = int(parts[i])
                title = re.sub(r"<[^>]+>", "", parts[i + 1]).strip()
                content = re.sub(r"<[^>]+>", " ", parts[i + 2] if i + 2 < len(parts) else "").strip()
                sections.append({"title": title, "content": content, "level": level})
    else:
        sections.append({"title": "", "content": clean_text, "level": 0})

    return {
        "sections": sections,
        "full_text": clean_text,
        "tables": tables,
        "format": "html",
        "parser": "regex",
        "section_count": len(sections),
        "success": True,
    }


def _parse_plain_text(text: str, preserve_structure: bool) -> dict[str, Any]:
    """플레인 텍스트 — 빈 줄 기반 단락 분할."""
    if preserve_structure:
        paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
        sections = [{"title": "", "content": p, "level": 0} for p in paragraphs]
    else:
        sections = [{"title": "", "content": text.strip(), "level": 0}]

    return {
        "sections": sections,
        "full_text": text,
        "tables": [],
        "format": "text",
        "parser": "regex",
        "section_count": len(sections),
        "success": True,
    }


def _parse_pdf(raw_bytes: bytes, preserve_structure: bool, extract_tables: bool) -> dict[str, Any]:
    """PDF 파싱 — PyMuPDF 또는 fallback."""
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=raw_bytes, filetype="pdf")
        sections = []
        full_text_parts = []
        tables = []

        for page_num, page in enumerate(doc):
            text = page.get_text("text")
            full_text_parts.append(text)

            if preserve_structure:
                sections.append({
                    "title": f"Page {page_num + 1}",
                    "content": text.strip(),
                    "page": page_num + 1,
                    "level": 0,
                })

            if extract_tables:
                try:
                    page_tables = page.find_tables()
                    for table in page_tables:
                        rows = table.extract()
                        if rows:
                            tables.append({"page": page_num + 1, "rows": rows})
                except Exception:
                    pass

        doc.close()
        full_text = "\n".join(full_text_parts)

        if not preserve_structure:
            sections = [{"title": "", "content": full_text.strip(), "level": 0}]

        return {
            "sections": sections,
            "full_text": full_text,
            "tables": tables,
            "format": "pdf",
            "parser": "pymupdf",
            "page_count": len(full_text_parts),
            "section_count": len(sections),
            "success": True,
        }

    except ImportError:
        logger.warning("PyMuPDF not installed, attempting basic PDF text extraction")
        return {
            "sections": [],
            "full_text": "",
            "tables": [],
            "format": "pdf",
            "success": False,
            "error": "PyMuPDF (fitz) not installed. Install with: pip install PyMuPDF",
        }
