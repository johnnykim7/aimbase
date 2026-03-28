"""CR-019: Document Reader Tools.

문서 파일을 구조화된 JSON/HTML로 변환하는 MCP Tool 모음.
- read_pptx: PPTX → 슬라이드별 JSON/HTML
- read_docx: DOCX → 단락/표/이미지 JSON/HTML
- read_excel: XLSX/CSV → 시트별 JSON
- read_pdf: PDF → 페이지별 텍스트/구조
"""

import base64
import csv
import io
import json
import logging
import os
import tempfile
from typing import Any

logger = logging.getLogger(__name__)


def read_pptx(file_base64: str, output_mode: str = "json") -> dict[str, Any]:
    """PPTX 파일 전체를 구조화된 JSON/HTML로 일괄 변환.

    Args:
        file_base64: base64 인코딩된 PPTX 파일
        output_mode: json 또는 html

    Returns:
        {slides: [...], slide_count, success}
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return {"success": False, "error": "python-pptx not installed"}

    try:
        file_bytes = base64.b64decode(file_base64)
        prs = Presentation(io.BytesIO(file_bytes))

        slides = []
        for idx, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": idx + 1,
                "layout": slide.slide_layout.name if slide.slide_layout else "",
                "elements": [],
            }

            for shape in slide.shapes:
                element = _extract_pptx_shape(shape)
                if element:
                    slide_data["elements"].append(element)

            # 표 처리
            if shape.has_table:
                table_data = _extract_table_from_pptx(shape.table)
                slide_data["elements"].append({"type": "table", "data": table_data})

            slides.append(slide_data)

        if output_mode == "html":
            html = _slides_to_html(slides)
            return {"html": html, "slide_count": len(slides), "success": True}

        return {"slides": slides, "slide_count": len(slides), "success": True}

    except Exception as e:
        logger.error("read_pptx failed: %s", e)
        return {"success": False, "error": str(e)}


def _extract_pptx_shape(shape) -> dict | None:
    """PPTX shape에서 데이터 추출."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    result = {
        "name": shape.name,
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }

    if shape.has_text_frame:
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs)
            if text.strip():
                paragraphs.append({
                    "text": text,
                    "level": para.level,
                    "bold": any(run.font.bold for run in para.runs if run.font.bold),
                })
        if paragraphs:
            result["type"] = "text"
            result["paragraphs"] = paragraphs
            return result

    if shape.has_table:
        result["type"] = "table"
        result["data"] = _extract_table_from_pptx(shape.table)
        return result

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            image = shape.image
            result["type"] = "image"
            result["content_type"] = image.content_type
            result["image_base64"] = base64.b64encode(image.blob).decode("utf-8")
            return result
        except Exception:
            result["type"] = "image"
            result["error"] = "Could not extract image"
            return result

    return None


def _extract_table_from_pptx(table) -> list[list[str]]:
    """PPTX 테이블 → 2D 배열."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip())
        rows.append(cells)
    return rows


def _slides_to_html(slides: list[dict]) -> str:
    """슬라이드 데이터 → HTML 변환."""
    parts = ["<!DOCTYPE html><html><head><meta charset='utf-8'><style>",
             ".slide{border:1px solid #ccc;margin:20px;padding:20px;page-break-after:always;}",
             ".slide h2{color:#333;border-bottom:2px solid #007bff;padding-bottom:8px;}",
             "table{border-collapse:collapse;margin:10px 0;}",
             "td,th{border:1px solid #ddd;padding:6px 10px;}",
             "img{max-width:400px;margin:10px 0;}",
             "</style></head><body>"]

    for s in slides:
        parts.append(f"<div class='slide'><h2>Slide {s['slide_number']}</h2>")
        for el in s.get("elements", []):
            if el.get("type") == "text":
                for p in el.get("paragraphs", []):
                    tag = "h3" if p.get("bold") or p.get("level", 0) == 0 and len(p["text"]) < 80 else "p"
                    parts.append(f"<{tag}>{_html_escape(p['text'])}</{tag}>")
            elif el.get("type") == "table":
                parts.append("<table>")
                for i, row in enumerate(el.get("data", [])):
                    tag = "th" if i == 0 else "td"
                    parts.append("<tr>" + "".join(f"<{tag}>{_html_escape(c)}</{tag}>" for c in row) + "</tr>")
                parts.append("</table>")
            elif el.get("type") == "image" and el.get("image_base64"):
                ct = el.get("content_type", "image/png")
                parts.append(f"<img src='data:{ct};base64,{el['image_base64']}' />")
        parts.append("</div>")

    parts.append("</body></html>")
    return "".join(parts)


def _html_escape(text: str) -> str:
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


# ── DOCX ──────────────────────────────────────────────────────

def read_docx(file_base64: str, output_mode: str = "json") -> dict[str, Any]:
    """DOCX 파일을 구조화된 JSON/HTML로 변환.

    Args:
        file_base64: base64 인코딩된 DOCX 파일
        output_mode: json 또는 html

    Returns:
        {paragraphs: [...], tables: [...], images: [...], success}
    """
    try:
        from docx import Document
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
    except ImportError:
        return {"success": False, "error": "python-docx not installed"}

    try:
        file_bytes = base64.b64decode(file_base64)
        doc = Document(io.BytesIO(file_bytes))

        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            paragraphs.append({
                "text": text,
                "style": para.style.name if para.style else "",
                "level": _docx_heading_level(para),
            })

        tables = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            tables.append(rows)

        # 이미지 추출
        images = []
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                try:
                    img_part = rel.target_part
                    images.append({
                        "content_type": img_part.content_type,
                        "image_base64": base64.b64encode(img_part.blob).decode("utf-8"),
                    })
                except Exception:
                    pass

        # 헤더/푸터
        headers_footers = []
        for section in doc.sections:
            for hf in [section.header, section.footer]:
                if hf is not None:
                    text = "\n".join(p.text.strip() for p in hf.paragraphs if p.text.strip())
                    if text:
                        headers_footers.append(text)

        result = {
            "paragraphs": paragraphs,
            "tables": tables,
            "images": images,
            "headers_footers": headers_footers,
            "paragraph_count": len(paragraphs),
            "table_count": len(tables),
            "image_count": len(images),
            "success": True,
        }

        if output_mode == "html":
            result["html"] = _docx_to_html(paragraphs, tables, images)

        return result

    except Exception as e:
        logger.error("read_docx failed: %s", e)
        return {"success": False, "error": str(e)}


def _docx_heading_level(para) -> int:
    """DOCX paragraph의 heading level 반환. 본문이면 -1."""
    style_name = para.style.name if para.style else ""
    if style_name.startswith("Heading"):
        try:
            return int(style_name.replace("Heading ", "").strip())
        except ValueError:
            pass
    return -1


def _docx_to_html(paragraphs, tables, images) -> str:
    """DOCX 데이터 → HTML."""
    parts = ["<!DOCTYPE html><html><head><meta charset='utf-8'><style>",
             "body{font-family:sans-serif;max-width:800px;margin:0 auto;padding:20px;}",
             "table{border-collapse:collapse;margin:10px 0;width:100%;}",
             "td,th{border:1px solid #ddd;padding:6px 10px;}",
             "img{max-width:600px;margin:10px 0;}",
             "</style></head><body>"]

    for p in paragraphs:
        level = p.get("level", -1)
        if level > 0:
            parts.append(f"<h{level}>{_html_escape(p['text'])}</h{level}>")
        else:
            parts.append(f"<p>{_html_escape(p['text'])}</p>")

    for table in tables:
        parts.append("<table>")
        for i, row in enumerate(table):
            tag = "th" if i == 0 else "td"
            parts.append("<tr>" + "".join(f"<{tag}>{_html_escape(c)}</{tag}>" for c in row) + "</tr>")
        parts.append("</table>")

    for img in images:
        ct = img.get("content_type", "image/png")
        parts.append(f"<img src='data:{ct};base64,{img['image_base64']}' />")

    parts.append("</body></html>")
    return "".join(parts)


# ── Excel ─────────────────────────────────────────────────────

def read_excel(file_base64: str, sheet_name: str = "", max_rows: int = 0) -> dict[str, Any]:
    """XLSX/CSV 파일을 시트별 JSON 데이터로 변환.

    Args:
        file_base64: base64 인코딩된 XLSX 또는 CSV 파일
        sheet_name: 특정 시트만 읽기 (빈 문자열이면 전체)
        max_rows: 최대 행 수 (0이면 전체)

    Returns:
        {sheets: [{name, headers, rows, row_count}], success}
    """
    try:
        file_bytes = base64.b64decode(file_base64)
    except Exception as e:
        return {"success": False, "error": f"Invalid base64: {e}"}

    # CSV 감지 (간단한 휴리스틱)
    try:
        text = file_bytes.decode("utf-8")
        if "," in text.split("\n")[0] and not text.startswith("PK"):
            return _read_csv(text, max_rows)
    except UnicodeDecodeError:
        pass

    # XLSX
    try:
        import openpyxl
    except ImportError:
        return {"success": False, "error": "openpyxl not installed"}

    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        sheets = []

        target_sheets = [sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.sheetnames

        for sname in target_sheets:
            ws = wb[sname]
            rows = []
            headers = []

            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i == 0:
                    headers = [str(c) if c is not None else "" for c in row]
                    continue
                if max_rows > 0 and i > max_rows:
                    break
                rows.append([_cell_to_str(c) for c in row])

            sheets.append({
                "name": sname,
                "headers": headers,
                "rows": rows,
                "row_count": len(rows),
            })

        wb.close()

        return {
            "sheets": sheets,
            "sheet_count": len(sheets),
            "success": True,
        }

    except Exception as e:
        logger.error("read_excel failed: %s", e)
        return {"success": False, "error": str(e)}


def _read_csv(text: str, max_rows: int) -> dict[str, Any]:
    """CSV 텍스트 → JSON."""
    reader = csv.reader(io.StringIO(text))
    headers = []
    rows = []

    for i, row in enumerate(reader):
        if i == 0:
            headers = row
            continue
        if max_rows > 0 and i > max_rows:
            break
        rows.append(row)

    return {
        "sheets": [{
            "name": "csv",
            "headers": headers,
            "rows": rows,
            "row_count": len(rows),
        }],
        "sheet_count": 1,
        "success": True,
    }


def _cell_to_str(cell) -> str:
    """Excel 셀 값을 문자열로 변환."""
    if cell is None:
        return ""
    if isinstance(cell, float) and cell == int(cell):
        return str(int(cell))
    return str(cell)


# ── PDF ───────────────────────────────────────────────────────

def read_pdf(
    file_base64: str,
    extract_images: bool = False,
    ocr_enabled: bool = False,
) -> dict[str, Any]:
    """PDF 파일에서 텍스트/구조를 추출.

    Args:
        file_base64: base64 인코딩된 PDF
        extract_images: 이미지 추출 여부
        ocr_enabled: OCR 사용 여부 (pytesseract 필요)

    Returns:
        {pages: [{page_number, text, tables}], page_count, success}
    """
    try:
        file_bytes = base64.b64decode(file_base64)
    except Exception as e:
        return {"success": False, "error": f"Invalid base64: {e}"}

    # pdfplumber 시도
    try:
        import pdfplumber
        return _read_pdf_pdfplumber(file_bytes, extract_images, ocr_enabled)
    except ImportError:
        pass

    # PyMuPDF fallback
    try:
        import fitz
        try:
            return _read_pdf_pymupdf(file_bytes, extract_images)
        except Exception as e:
            return {"success": False, "error": f"PDF parsing failed: {e}"}
    except ImportError:
        pass

    return {"success": False, "error": "No PDF library available. Install pdfplumber or PyMuPDF."}


def _read_pdf_pdfplumber(file_bytes: bytes, extract_images: bool, ocr_enabled: bool) -> dict[str, Any]:
    """pdfplumber로 PDF 읽기."""
    import pdfplumber

    pages = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages):
            page_data = {
                "page_number": i + 1,
                "text": page.extract_text() or "",
                "width": page.width,
                "height": page.height,
            }

            # 테이블 추출
            tables = page.extract_tables()
            if tables:
                page_data["tables"] = tables

            # 이미지 추출
            if extract_images and page.images:
                page_data["image_count"] = len(page.images)

            pages.append(page_data)

    total_text = sum(len(p["text"]) for p in pages)

    return {
        "pages": pages,
        "page_count": len(pages),
        "total_characters": total_text,
        "success": True,
    }


def _read_pdf_pymupdf(file_bytes: bytes, extract_images: bool) -> dict[str, Any]:
    """PyMuPDF(fitz)로 PDF 읽기."""
    import fitz

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = []

    for i, page in enumerate(doc):
        page_data = {
            "page_number": i + 1,
            "text": page.get_text(),
            "width": page.rect.width,
            "height": page.rect.height,
        }

        if extract_images:
            images = []
            for img_info in page.get_images(full=True):
                try:
                    xref = img_info[0]
                    img = doc.extract_image(xref)
                    images.append({
                        "content_type": f"image/{img['ext']}",
                        "image_base64": base64.b64encode(img["image"]).decode("utf-8"),
                    })
                except Exception:
                    pass
            if images:
                page_data["images"] = images

        pages.append(page_data)

    doc.close()
    total_text = sum(len(p["text"]) for p in pages)

    return {
        "pages": pages,
        "page_count": len(pages),
        "total_characters": total_text,
        "success": True,
    }
