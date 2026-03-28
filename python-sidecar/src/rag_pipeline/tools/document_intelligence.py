"""CR-013: Document Intelligence — 스키마 기반 문서 생성 (Level 3).

LLM이 구조화된 JSON 스키마를 생성하면, 렌더러가 PPTX/DOCX/PDF로 변환.
CR-018(코드 실행 기반)과 달리 테마/브랜드/후처리를 렌더러가 제어.

지원:
- Presentation (PPTX): 슬라이드 레이아웃 6종
- Document (DOCX): 헤딩/본문/표/목록/이미지
- PDF: reportlab 기반 문서
"""

import base64
import copy
import json
import logging
import os
import tempfile
from typing import Any

logger = logging.getLogger(__name__)

# ─── 테마 프리셋 ───────────────────────────────────────────────

THEMES: dict[str, dict[str, Any]] = {
    "default": {
        "name": "기본",
        "primary_color": "#1a73e8",
        "secondary_color": "#34a853",
        "accent_color": "#ea4335",
        "background_color": "#ffffff",
        "text_color": "#202124",
        "font_family": "Malgun Gothic",
        "font_family_fallback": "Arial",
        "font_size_title": 28,
        "font_size_subtitle": 18,
        "font_size_heading": 22,
        "font_size_body": 12,
        "font_size_caption": 9,
        "margin_cm": 2.54,
        "line_spacing": 1.5,
    },
    "corporate_blue": {
        "name": "기업 블루",
        "primary_color": "#003366",
        "secondary_color": "#336699",
        "accent_color": "#ff9900",
        "background_color": "#ffffff",
        "text_color": "#333333",
        "font_family": "Malgun Gothic",
        "font_family_fallback": "Calibri",
        "font_size_title": 32,
        "font_size_subtitle": 20,
        "font_size_heading": 24,
        "font_size_body": 12,
        "font_size_caption": 9,
        "margin_cm": 2.54,
        "line_spacing": 1.5,
    },
    "modern_dark": {
        "name": "모던 다크",
        "primary_color": "#bb86fc",
        "secondary_color": "#03dac6",
        "accent_color": "#cf6679",
        "background_color": "#121212",
        "text_color": "#e0e0e0",
        "font_family": "Malgun Gothic",
        "font_family_fallback": "Segoe UI",
        "font_size_title": 30,
        "font_size_subtitle": 18,
        "font_size_heading": 24,
        "font_size_body": 12,
        "font_size_caption": 9,
        "margin_cm": 2.54,
        "line_spacing": 1.5,
    },
    "minimal": {
        "name": "미니멀",
        "primary_color": "#000000",
        "secondary_color": "#666666",
        "accent_color": "#0066cc",
        "background_color": "#ffffff",
        "text_color": "#333333",
        "font_family": "Malgun Gothic",
        "font_family_fallback": "Helvetica",
        "font_size_title": 36,
        "font_size_subtitle": 16,
        "font_size_heading": 22,
        "font_size_body": 11,
        "font_size_caption": 8,
        "margin_cm": 3.0,
        "line_spacing": 1.6,
    },
    "warm": {
        "name": "따뜻한 톤",
        "primary_color": "#c0392b",
        "secondary_color": "#e67e22",
        "accent_color": "#f39c12",
        "background_color": "#fdf6ec",
        "text_color": "#2c3e50",
        "font_family": "Malgun Gothic",
        "font_family_fallback": "Georgia",
        "font_size_title": 30,
        "font_size_subtitle": 18,
        "font_size_heading": 24,
        "font_size_body": 12,
        "font_size_caption": 9,
        "margin_cm": 2.54,
        "line_spacing": 1.5,
    },
}

# ─── 스키마 검증 ───────────────────────────────────────────────

VALID_DOC_TYPES = {"presentation", "document"}
VALID_SLIDE_LAYOUTS = {"title", "title_content", "two_column", "image_text", "chart", "table", "blank"}
VALID_CONTENT_TYPES = {"text", "bullets", "numbered", "image", "table", "chart"}
VALID_SECTION_TYPES = {"heading", "paragraph", "bullets", "numbered", "table", "image", "page_break"}
VALID_CHART_TYPES = {"bar", "line", "pie", "column", "area"}


def validate_schema(schema: dict[str, Any]) -> dict[str, Any]:
    """문서 스키마 검증.

    Args:
        schema: 문서 스키마 JSON

    Returns:
        {valid: bool, errors: [...], warnings: [...]}
    """
    errors: list[str] = []
    warnings: list[str] = []

    doc_type = schema.get("type")
    if not doc_type:
        errors.append("'type' 필드 필수 (presentation 또는 document)")
        return {"valid": False, "errors": errors, "warnings": warnings}

    if doc_type not in VALID_DOC_TYPES:
        errors.append(f"type '{doc_type}'은 지원하지 않음. 가능: {VALID_DOC_TYPES}")

    if doc_type == "presentation":
        _validate_presentation(schema, errors, warnings)
    elif doc_type == "document":
        _validate_document(schema, errors, warnings)

    return {"valid": len(errors) == 0, "errors": errors, "warnings": warnings}


def _validate_presentation(schema: dict, errors: list, warnings: list):
    slides = schema.get("slides", [])
    if not slides:
        errors.append("slides 배열이 비어있음")
        return

    for i, slide in enumerate(slides):
        layout = slide.get("layout", "")
        if layout not in VALID_SLIDE_LAYOUTS:
            errors.append(f"slide[{i}]: layout '{layout}' 미지원. 가능: {VALID_SLIDE_LAYOUTS}")

        if layout in ("title", "title_content") and not slide.get("title"):
            warnings.append(f"slide[{i}]: '{layout}' 레이아웃에 title 없음")

        if layout == "chart":
            chart = slide.get("chart", {})
            if chart.get("type") not in VALID_CHART_TYPES:
                errors.append(f"slide[{i}]: chart type '{chart.get('type')}' 미지원")

        if layout == "table":
            table = slide.get("table", {})
            if not table.get("headers"):
                errors.append(f"slide[{i}]: table에 headers 필수")

    if len(slides) > 50:
        warnings.append(f"슬라이드 {len(slides)}개 — 생성 시간이 길어질 수 있음")


def _validate_document(schema: dict, errors: list, warnings: list):
    sections = schema.get("sections", [])
    if not sections:
        errors.append("sections 배열이 비어있음")
        return

    for i, sec in enumerate(sections):
        sec_type = sec.get("type", "")
        if sec_type not in VALID_SECTION_TYPES:
            errors.append(f"section[{i}]: type '{sec_type}' 미지원. 가능: {VALID_SECTION_TYPES}")

        if sec_type == "heading" and not sec.get("text"):
            warnings.append(f"section[{i}]: heading에 text 없음")

        if sec_type == "table" and not sec.get("headers"):
            errors.append(f"section[{i}]: table에 headers 필수")


# ─── 후처리 규칙 ───────────────────────────────────────────────

MAX_BULLET_CHARS = 120
MAX_SLIDE_BULLETS = 7
MAX_TABLE_ROWS_PER_SLIDE = 12


def _postprocess_text(text: str, max_chars: int = MAX_BULLET_CHARS) -> str:
    """텍스트 길이 초과 시 축약."""
    if len(text) <= max_chars:
        return text
    # 마지막 완전한 단어까지 자름
    truncated = text[:max_chars - 3].rsplit(" ", 1)[0]
    return truncated + "..."


def _postprocess_bullets(items: list[str]) -> list[str]:
    """불릿 아이템 후처리: 개수 제한 + 텍스트 축약."""
    processed = [_postprocess_text(item) for item in items]
    if len(processed) > MAX_SLIDE_BULLETS:
        processed = processed[:MAX_SLIDE_BULLETS - 1]
        remaining = len(items) - MAX_SLIDE_BULLETS + 1
        processed.append(f"외 {remaining}개 항목...")
    return processed


def _postprocess_table_rows(rows: list[list[str]], max_rows: int = MAX_TABLE_ROWS_PER_SLIDE) -> list[list[list[str]]]:
    """테이블 행 수 초과 시 분할."""
    if len(rows) <= max_rows:
        return [rows]
    chunks = []
    for i in range(0, len(rows), max_rows):
        chunks.append(rows[i:i + max_rows])
    return chunks


# ─── 색상 유틸 ─────────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """#RRGGBB → (R, G, B)."""
    h = hex_color.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


# ─── Presentation 렌더러 (PPTX) ───────────────────────────────

def _render_presentation(schema: dict[str, Any], theme: dict[str, Any]) -> bytes:
    """스키마 → PPTX 바이트."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    primary_rgb = RGBColor(*_hex_to_rgb(theme["primary_color"]))
    secondary_rgb = RGBColor(*_hex_to_rgb(theme["secondary_color"]))
    text_rgb = RGBColor(*_hex_to_rgb(theme["text_color"]))
    bg_rgb = RGBColor(*_hex_to_rgb(theme["background_color"]))
    accent_rgb = RGBColor(*_hex_to_rgb(theme.get("accent_color", "#ea4335")))

    font_name = theme.get("font_family", "Malgun Gothic")
    title_size = Pt(theme.get("font_size_title", 28))
    subtitle_size = Pt(theme.get("font_size_subtitle", 18))
    heading_size = Pt(theme.get("font_size_heading", 22))
    body_size = Pt(theme.get("font_size_body", 12))

    for slide_data in schema.get("slides", []):
        layout = slide_data.get("layout", "blank")
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        # 배경색 설정
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_rgb

        if layout == "title":
            _add_title_slide(slide, slide_data, font_name, title_size, subtitle_size, primary_rgb, text_rgb)
        elif layout == "title_content":
            _add_title_content_slide(slide, slide_data, font_name, heading_size, body_size, primary_rgb, text_rgb)
        elif layout == "two_column":
            _add_two_column_slide(slide, slide_data, font_name, heading_size, body_size, primary_rgb, text_rgb)
        elif layout == "table":
            _add_table_slide(slide, slide_data, font_name, heading_size, body_size, primary_rgb, text_rgb, bg_rgb)
        elif layout == "chart":
            _add_chart_slide(slide, slide_data, font_name, heading_size, primary_rgb, text_rgb)
        elif layout == "image_text":
            _add_image_text_slide(slide, slide_data, font_name, heading_size, body_size, primary_rgb, text_rgb)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        tmp.seek(0)
        data = open(tmp.name, "rb").read()
    os.unlink(tmp.name)
    return data


def _set_font(run, font_name, size, color, bold=False):
    run.font.name = font_name
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold


def _add_title_slide(slide, data, font_name, title_size, subtitle_size, primary_rgb, text_rgb):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    # 제목
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = data.get("title", "")
    _set_font(run, font_name, title_size, primary_rgb, bold=True)

    # 부제목
    subtitle = data.get("subtitle", "")
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(2), Inches(4.0), Inches(9.3), Inches(1))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        _set_font(run2, font_name, subtitle_size, text_rgb)


def _add_title_content_slide(slide, data, font_name, heading_size, body_size, primary_rgb, text_rgb):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    # 제목
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = data.get("title", "")
    _set_font(run, font_name, heading_size, primary_rgb, bold=True)

    # 구분선
    from pptx.util import Emu
    line = slide.shapes.add_connector(1, Inches(0.7), Inches(1.3), Inches(11.9), Inches(1.3))
    line.line.color.rgb = primary_rgb
    line.line.width = Emu(18000)

    # 콘텐츠
    content = data.get("content", {})
    _render_content_block(slide, content, font_name, body_size, text_rgb,
                          left=Inches(0.7), top=Inches(1.5), width=Inches(11.9), height=Inches(5.5))


def _add_two_column_slide(slide, data, font_name, heading_size, body_size, primary_rgb, text_rgb):
    from pptx.util import Inches

    # 제목
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = data.get("title", "")
    _set_font(run, font_name, heading_size, primary_rgb, bold=True)

    # 좌측
    left_content = data.get("left", {})
    _render_content_block(slide, left_content, font_name, body_size, text_rgb,
                          left=Inches(0.7), top=Inches(1.5), width=Inches(5.7), height=Inches(5.5))

    # 우측
    right_content = data.get("right", {})
    _render_content_block(slide, right_content, font_name, body_size, text_rgb,
                          left=Inches(6.9), top=Inches(1.5), width=Inches(5.7), height=Inches(5.5))


def _add_table_slide(slide, data, font_name, heading_size, body_size, primary_rgb, text_rgb, bg_rgb):
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    # 제목
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = data.get("title", "")
    _set_font(run, font_name, heading_size, primary_rgb, bold=True)

    table_data = data.get("table", {})
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    row_chunks = _postprocess_table_rows(rows)

    # 첫 번째 청크만 현재 슬라이드에 표시
    chunk = row_chunks[0] if row_chunks else []
    cols = len(headers) if headers else (len(chunk[0]) if chunk else 1)
    total_rows = 1 + len(chunk)  # header + data

    tbl_shape = slide.shapes.add_table(total_rows, cols, Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.5))
    tbl = tbl_shape.table

    # 헤더
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = str(h)
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                _set_font(run, font_name, body_size, RGBColor(255, 255, 255), bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = primary_rgb

    # 데이터
    for ri, row in enumerate(chunk):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = _postprocess_text(str(val), 80)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _set_font(run, font_name, body_size, text_rgb)


def _add_chart_slide(slide, data, font_name, heading_size, primary_rgb, text_rgb):
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    # 제목
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = data.get("title", "")
    _set_font(run, font_name, heading_size, primary_rgb, bold=True)

    chart_info = data.get("chart", {})
    chart_type_str = chart_info.get("type", "bar")
    categories = chart_info.get("categories", [])
    series_list = chart_info.get("series", [])

    type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "area": XL_CHART_TYPE.AREA,
    }
    xl_type = type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = categories

    for s in series_list:
        chart_data.add_series(s.get("name", ""), s.get("values", []))

    slide.shapes.add_chart(xl_type, Inches(1), Inches(1.5), Inches(11.3), Inches(5.5), chart_data)


def _add_image_text_slide(slide, data, font_name, heading_size, body_size, primary_rgb, text_rgb):
    from pptx.util import Inches

    # 제목
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = data.get("title", "")
    _set_font(run, font_name, heading_size, primary_rgb, bold=True)

    # 텍스트 (우측)
    text_content = data.get("content", data.get("text", ""))
    if isinstance(text_content, dict):
        _render_content_block(slide, text_content, font_name, body_size, text_rgb,
                              left=Inches(6.9), top=Inches(1.5), width=Inches(5.7), height=Inches(5.5))
    else:
        txBox2 = slide.shapes.add_textbox(Inches(6.9), Inches(1.5), Inches(5.7), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = str(text_content)
        _set_font(run2, font_name, body_size, text_rgb)

    # 이미지 (좌측) — base64 또는 placeholder
    image_data = data.get("image", {})
    image_b64 = image_data.get("base64", "")
    if image_b64:
        import io
        img_bytes = base64.b64decode(image_b64)
        img_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(img_stream, Inches(0.7), Inches(1.5), Inches(5.7), Inches(5.5))
    else:
        # placeholder 텍스트
        txBox3 = slide.shapes.add_textbox(Inches(0.7), Inches(3), Inches(5.7), Inches(1))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        run3 = p3.add_run()
        run3.text = image_data.get("alt", "[이미지]")
        _set_font(run3, font_name, body_size, text_rgb)


def _render_content_block(slide, content, font_name, body_size, text_rgb, left, top, width, height):
    """콘텐츠 블록 렌더링 (bullets, text, numbered)."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    if not content:
        return

    content_type = content.get("type", "text") if isinstance(content, dict) else "text"

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if content_type == "bullets":
        items = content.get("items", [])
        items = _postprocess_bullets(items)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"• {item}"
            _set_font(run, font_name, body_size, text_rgb)
    elif content_type == "numbered":
        items = content.get("items", [])
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            run = p.add_run()
            run.text = f"{i + 1}. {_postprocess_text(item)}"
            _set_font(run, font_name, body_size, text_rgb)
    else:
        text = content.get("text", str(content)) if isinstance(content, dict) else str(content)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        _set_font(run, font_name, body_size, text_rgb)


# ─── Document 렌더러 (DOCX) ───────────────────────────────────

def _render_document(schema: dict[str, Any], theme: dict[str, Any]) -> bytes:
    """스키마 → DOCX 바이트."""
    from docx import Document
    from docx.shared import Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # 여백 설정
    for section in doc.sections:
        margin = Cm(theme.get("margin_cm", 2.54))
        section.top_margin = margin
        section.bottom_margin = margin
        section.left_margin = margin
        section.right_margin = margin

    font_name = theme.get("font_family", "Malgun Gothic")
    primary_rgb = RGBColor(*_hex_to_rgb(theme["primary_color"]))
    text_rgb = RGBColor(*_hex_to_rgb(theme["text_color"]))

    # 제목 (스키마 최상위)
    doc_title = schema.get("title", "")
    if doc_title:
        p = doc.add_heading(doc_title, level=0)
        for run in p.runs:
            run.font.name = font_name
            run.font.color.rgb = primary_rgb

    for sec in schema.get("sections", []):
        sec_type = sec.get("type", "paragraph")

        if sec_type == "heading":
            level = sec.get("level", 1)
            p = doc.add_heading(sec.get("text", ""), level=min(level, 4))
            for run in p.runs:
                run.font.name = font_name
                run.font.color.rgb = primary_rgb

        elif sec_type == "paragraph":
            text = sec.get("text", "")
            p = doc.add_paragraph(text)
            p.paragraph_format.line_spacing = theme.get("line_spacing", 1.5)
            for run in p.runs:
                run.font.name = font_name
                run.font.size = Pt(theme.get("font_size_body", 12))
                run.font.color.rgb = text_rgb
            if sec.get("bold"):
                for run in p.runs:
                    run.bold = True

        elif sec_type == "bullets":
            items = sec.get("items", [])
            for item in items:
                p = doc.add_paragraph(str(item), style="List Bullet")
                for run in p.runs:
                    run.font.name = font_name
                    run.font.size = Pt(theme.get("font_size_body", 12))

        elif sec_type == "numbered":
            items = sec.get("items", [])
            for item in items:
                p = doc.add_paragraph(str(item), style="List Number")
                for run in p.runs:
                    run.font.name = font_name
                    run.font.size = Pt(theme.get("font_size_body", 12))

        elif sec_type == "table":
            headers = sec.get("headers", [])
            rows = sec.get("rows", [])
            cols = len(headers) if headers else (len(rows[0]) if rows else 1)
            table = doc.add_table(rows=1, cols=cols)
            table.style = "Table Grid"

            # 헤더
            for ci, h in enumerate(headers):
                cell = table.rows[0].cells[ci]
                cell.text = str(h)
                for p in cell.paragraphs:
                    for run in p.runs:
                        run.font.name = font_name
                        run.font.bold = True
                        run.font.size = Pt(theme.get("font_size_body", 12))

            # 데이터
            for row_data in rows:
                row = table.add_row()
                for ci, val in enumerate(row_data):
                    row.cells[ci].text = str(val)
                    for p in row.cells[ci].paragraphs:
                        for run in p.runs:
                            run.font.name = font_name
                            run.font.size = Pt(theme.get("font_size_body", 12))

        elif sec_type == "image":
            img_b64 = sec.get("base64", "")
            if img_b64:
                import io
                img_bytes = base64.b64decode(img_b64)
                img_stream = io.BytesIO(img_bytes)
                width_cm = sec.get("width_cm", 15)
                doc.add_picture(img_stream, width=Cm(width_cm))
            else:
                doc.add_paragraph(f"[이미지: {sec.get('alt', '')}]")

        elif sec_type == "page_break":
            doc.add_page_break()

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
        doc.save(tmp.name)
        data = open(tmp.name, "rb").read()
    os.unlink(tmp.name)
    return data


# ─── PDF 렌더러 ────────────────────────────────────────────────

def _render_pdf(schema: dict[str, Any], theme: dict[str, Any]) -> bytes:
    """스키마 → PDF 바이트."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib.colors import HexColor
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.enums import TA_LEFT, TA_CENTER

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        tmp_path = tmp.name

    margin = theme.get("margin_cm", 2.54) * cm
    doc = SimpleDocTemplate(tmp_path, pagesize=A4,
                            topMargin=margin, bottomMargin=margin,
                            leftMargin=margin, rightMargin=margin)

    styles = getSampleStyleSheet()
    font_name = "Helvetica"  # reportlab 기본 폰트 (한글은 별도 등록 필요)
    primary = HexColor(theme["primary_color"])
    text_color = HexColor(theme["text_color"])

    title_style = ParagraphStyle("DocTitle", parent=styles["Title"],
                                 fontSize=theme.get("font_size_title", 28),
                                 textColor=primary, fontName=font_name)
    heading_style = ParagraphStyle("DocHeading", parent=styles["Heading1"],
                                   fontSize=theme.get("font_size_heading", 22),
                                   textColor=primary, fontName=font_name)
    body_style = ParagraphStyle("DocBody", parent=styles["Normal"],
                                fontSize=theme.get("font_size_body", 12),
                                textColor=text_color, fontName=font_name,
                                leading=theme.get("font_size_body", 12) * theme.get("line_spacing", 1.5))

    elements: list = []

    doc_title = schema.get("title", "")
    if doc_title:
        elements.append(Paragraph(doc_title, title_style))
        elements.append(Spacer(1, 0.5 * cm))

    for sec in schema.get("sections", []):
        sec_type = sec.get("type", "paragraph")

        if sec_type == "heading":
            level = sec.get("level", 1)
            h_style = ParagraphStyle(f"H{level}", parent=heading_style,
                                     fontSize=max(22 - (level - 1) * 3, 12))
            elements.append(Paragraph(sec.get("text", ""), h_style))
            elements.append(Spacer(1, 0.3 * cm))

        elif sec_type == "paragraph":
            elements.append(Paragraph(sec.get("text", ""), body_style))
            elements.append(Spacer(1, 0.2 * cm))

        elif sec_type in ("bullets", "numbered"):
            items = sec.get("items", [])
            for i, item in enumerate(items):
                prefix = f"{i + 1}. " if sec_type == "numbered" else "• "
                elements.append(Paragraph(f"{prefix}{item}", body_style))
            elements.append(Spacer(1, 0.2 * cm))

        elif sec_type == "table":
            headers = sec.get("headers", [])
            rows = sec.get("rows", [])
            table_data = [headers] + rows if headers else rows

            if table_data:
                t = Table(table_data)
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), primary),
                    ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#ffffff")),
                    ("FONTSIZE", (0, 0), (-1, -1), theme.get("font_size_body", 12)),
                    ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#cccccc")),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]))
                elements.append(t)
                elements.append(Spacer(1, 0.3 * cm))

        elif sec_type == "page_break":
            elements.append(PageBreak())

    doc.build(elements)
    data = open(tmp_path, "rb").read()
    os.unlink(tmp_path)
    return data


# ─── 메인 API ──────────────────────────────────────────────────

def generate_from_schema(
    schema: str | dict,
    output_format: str = "pptx",
    theme_name: str = "default",
    custom_theme: str = "{}",
    output_filename: str = "",
) -> dict[str, Any]:
    """구조화된 JSON 스키마로 문서 생성 (CR-013 Level 3).

    Args:
        schema: 문서 스키마 (JSON 문자열 또는 dict)
        output_format: 출력 포맷 (pptx, docx, pdf)
        theme_name: 테마 프리셋명 (default, corporate_blue, modern_dark, minimal, warm)
        custom_theme: 커스텀 테마 오버라이드 JSON (프리셋 위에 병합)
        output_filename: 출력 파일명

    Returns:
        {file_base64, filename, format, size_bytes, success, schema_validation}
    """
    # 스키마 파싱
    if isinstance(schema, str):
        try:
            schema = json.loads(schema)
        except json.JSONDecodeError as e:
            return {"success": False, "error": f"Invalid schema JSON: {e}"}

    # 스키마 검증
    validation = validate_schema(schema)
    if not validation["valid"]:
        return {
            "success": False,
            "error": "Schema validation failed",
            "schema_validation": validation,
        }

    # 테마 선택 + 커스텀 오버라이드
    theme = copy.deepcopy(THEMES.get(theme_name, THEMES["default"]))
    if custom_theme and custom_theme != "{}":
        try:
            overrides = json.loads(custom_theme) if isinstance(custom_theme, str) else custom_theme
            theme.update(overrides)
        except (json.JSONDecodeError, TypeError):
            pass

    # 파일명
    fmt = output_format.lower().strip()
    ext_map = {"pptx": ".pptx", "docx": ".docx", "pdf": ".pdf"}
    if fmt not in ext_map:
        return {"success": False, "error": f"Schema rendering supports: pptx, docx, pdf. Got: {fmt}"}

    filename = output_filename or f"document{ext_map[fmt]}"
    if not filename.endswith(ext_map[fmt]):
        filename += ext_map[fmt]

    try:
        if fmt == "pptx":
            file_bytes = _render_presentation(schema, theme)
        elif fmt == "docx":
            file_bytes = _render_document(schema, theme)
        elif fmt == "pdf":
            file_bytes = _render_pdf(schema, theme)
        else:
            return {"success": False, "error": f"Unsupported format: {fmt}"}

        file_base64 = base64.b64encode(file_bytes).decode("utf-8")

        logger.info("Schema-based document generated: %s (%d bytes)", filename, len(file_bytes))
        return {
            "file_base64": file_base64,
            "filename": filename,
            "format": fmt,
            "size_bytes": len(file_bytes),
            "theme": theme_name,
            "success": True,
            "schema_validation": validation,
        }

    except Exception as e:
        logger.error("Schema-based document generation failed: %s", e)
        return {
            "success": False,
            "error": f"Rendering failed: {type(e).__name__}: {str(e)}",
            "schema_validation": validation,
        }


def list_themes() -> dict[str, Any]:
    """사용 가능한 테마 프리셋 목록 반환."""
    result = {}
    for key, theme in THEMES.items():
        result[key] = {
            "name": theme["name"],
            "primary_color": theme["primary_color"],
            "secondary_color": theme["secondary_color"],
            "background_color": theme["background_color"],
            "font_family": theme["font_family"],
        }
    return {"themes": result, "count": len(result), "success": True}
