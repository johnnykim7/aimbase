"""Ingest Document MCP Tool.

Orchestrates the full pipeline: parse → chunk → embed → store in pgvector.
"""

import logging
import os

from rag_pipeline.db import delete_embeddings_by_source, store_embeddings
from rag_pipeline.tools.chunker import chunk_document
from rag_pipeline.tools.embedder import embed_texts

logger = logging.getLogger(__name__)

BATCH_SIZE = 20


def _recommend_strategy(file_ext: str, explicit_strategy: str = "") -> str:
    """Auto-recommend chunking strategy based on file type."""
    if explicit_strategy and explicit_strategy != "auto":
        return explicit_strategy

    ext_strategy_map = {
        ".pptx": "fixed",
        ".xlsx": "fixed",
        ".csv": "fixed",
        ".pdf": "semantic",
        ".docx": "semantic",
        ".md": "recursive",
        ".html": "recursive",
        ".htm": "recursive",
        ".txt": "fixed",
    }
    return ext_strategy_map.get(file_ext.lower(), "fixed")


def ingest_document(
    source_id: str,
    content: str,
    document_id: str = "",
    chunking_strategy: str = "semantic",
    chunking_config: dict | None = None,
    embedding_model: str = "",
) -> dict:
    """Full ingestion pipeline: chunk → embed → store.

    Args:
        source_id: Knowledge source ID
        content: Raw document text
        document_id: Optional document identifier
        chunking_strategy: "semantic" | "recursive" | "fixed"
        chunking_config: Chunking parameters
        embedding_model: Embedding model name

    Returns:
        {source_id, document_id, chunks_created, success, errors}
    """
    errors: list[dict] = []
    doc_id = document_id or source_id
    # Strip NUL bytes — PostgreSQL text columns reject \x00
    content = content.replace("\x00", "")

    try:
        # 1) Delete existing embeddings for re-sync
        deleted = delete_embeddings_by_source(source_id)
        if deleted > 0:
            logger.info("Deleted %d existing embeddings for source '%s'", deleted, source_id)

        # 2) Chunk
        chunks = chunk_document(content, chunking_strategy, chunking_config)
        if not chunks:
            return {
                "source_id": source_id,
                "document_id": doc_id,
                "chunks_created": 0,
                "success": True,
                "errors": [],
            }

        logger.info("Created %d chunks for source '%s'", len(chunks), source_id)

        # 3) Embed + Store in batches
        total_stored = 0
        for i in range(0, len(chunks), BATCH_SIZE):
            batch = chunks[i : i + BATCH_SIZE]
            texts = [c["content"] for c in batch]

            try:
                result = embed_texts(texts, embedding_model)
                vectors = result["embeddings"]
                stored = store_embeddings(source_id, batch, vectors, doc_id)
                total_stored += stored
            except Exception as e:
                logger.warning("Batch %d failed: %s", i // BATCH_SIZE, str(e))
                errors.append({"batch": i // BATCH_SIZE, "error": str(e)})

        return {
            "source_id": source_id,
            "document_id": doc_id,
            "chunks_created": total_stored,
            "success": len(errors) == 0,
            "errors": errors,
        }

    except Exception as e:
        logger.error("Ingestion failed for source '%s': %s", source_id, str(e))
        return {
            "source_id": source_id,
            "document_id": doc_id,
            "chunks_created": 0,
            "success": False,
            "errors": [{"error": str(e)}],
        }


# ── File parsing helpers ────────────────────────────────────────────


def _parse_file_advanced(full_path: str) -> str:
    """Unstructured 라이브러리로 파싱을 시도하고, 실패하면 기존 파서로 fallback."""
    ext = os.path.splitext(full_path)[1].lower()

    # Unstructured가 효과적인 포맷: PDF, DOCX, PPTX, HTML
    if ext in (".pdf", ".docx", ".pptx", ".html", ".htm"):
        try:
            from rag_pipeline.tools.advanced_parser import _HAS_UNSTRUCTURED
            if _HAS_UNSTRUCTURED:
                import base64
                from rag_pipeline.tools.advanced_parser import advanced_parse
                with open(full_path, "rb") as f:
                    file_bytes = f.read()
                b64 = base64.b64encode(file_bytes).decode("ascii")
                ft = ext.lstrip(".")
                if ft == "htm":
                    ft = "html"
                result = advanced_parse(b64, file_type=ft, config={"extract_tables": True})
                if result.get("success") and result.get("full_text", "").strip():
                    parser_used = result.get("parser", "unknown")
                    logger.info("Advanced parsing (%s) succeeded for '%s': %d sections, %d tables",
                                parser_used, full_path, result.get("section_count", 0), result.get("table_count", 0))
                    return result["full_text"]
        except Exception as e:
            logger.warning("Advanced parsing failed for '%s', falling back: %s", full_path, e)

    return _parse_file(full_path)


def _parse_file(full_path: str) -> str:
    """Parse a file into plain text based on its extension.

    Supports txt/md/csv/html directly, pdf/docx/pptx via optional libraries.
    """
    ext = os.path.splitext(full_path)[1].lower()

    # Plain text formats
    if ext in (".txt", ".md", ".csv", ".html", ".htm", ".json", ".xml", ".yaml", ".yml"):
        with open(full_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    # PDF via PyMuPDF (pymupdf)
    if ext == ".pdf":
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(full_path)
            pages = []
            for page in doc:
                pages.append(page.get_text())
            doc.close()
            return "\n\n".join(pages)
        except ImportError:
            logger.warning("PyMuPDF not installed, trying basic PDF text extraction")
        except Exception as e:
            logger.warning("PyMuPDF failed for '%s': %s", full_path, e)
        # Fallback: read as binary and try to extract text
        return _fallback_read(full_path, ext)

    # DOCX via python-docx
    if ext == ".docx":
        try:
            from docx import Document

            doc = Document(full_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs)
        except ImportError:
            logger.warning("python-docx not installed for docx parsing")
        except Exception as e:
            logger.warning("python-docx failed for '%s': %s", full_path, e)
        return _fallback_read(full_path, ext)

    # PPTX via python-pptx
    if ext == ".pptx":
        try:
            from pptx import Presentation

            prs = Presentation(full_path)
            slides_text = []
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_parts = []
                for shape in slide.shapes:
                    if shape.has_table:
                        # 테이블 데이터 추출
                        table = shape.table
                        rows = []
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            rows.append(" | ".join(cells))
                        slide_parts.append("\n".join(rows))
                    elif hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text)
                if slide_parts:
                    slide_content = "\n".join(slide_parts)
                    slides_text.append(f"[슬라이드 {slide_idx}]\n{slide_content}")
            return "\n\n---\n\n".join(slides_text)
        except ImportError:
            logger.warning("python-pptx not installed for pptx parsing")
        except Exception as e:
            logger.warning("python-pptx failed for '%s': %s", full_path, e)
        return _fallback_read(full_path, ext)

    # Unknown extension: try reading as text
    return _fallback_read(full_path, ext)


def _fallback_read(full_path: str, ext: str) -> str:
    """Last-resort: read file as UTF-8 text, ignoring decode errors."""
    try:
        with open(full_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        if content.strip():
            return content
    except Exception:
        pass
    raise ValueError(f"Cannot parse file with extension '{ext}': {full_path}")


def ingest_file(
    source_id: str,
    file_path: str,
    storage_base_path: str,
    document_id: str = "",
    chunking_strategy: str = "semantic",
    chunking_config: dict | None = None,
    embedding_model: str = "",
) -> dict:
    """File-based ingestion: read file → parse → chunk → embed → store.

    The sidecar handles the entire pipeline, including file parsing.
    BE only needs to pass the file path and storage base path.

    Args:
        source_id: Knowledge source ID
        file_path: Relative file path within storage (e.g. tenant/category/file.pdf)
        storage_base_path: Absolute base path of storage directory
        document_id: Optional document identifier (defaults to file_path)
        chunking_strategy: "semantic" | "recursive" | "fixed"
        chunking_config: Chunking parameters
        embedding_model: Embedding model name

    Returns:
        {source_id, document_id, chunks_created, success, errors}
    """
    doc_id = document_id or file_path

    try:
        # 1) Resolve full path
        full_path = os.path.join(storage_base_path, file_path)
        full_path = os.path.abspath(full_path)

        if not os.path.isfile(full_path):
            return {
                "source_id": source_id,
                "document_id": doc_id,
                "chunks_created": 0,
                "success": False,
                "errors": [{"error": f"File not found: {full_path}"}],
            }

        logger.info("Parsing file '%s' for source '%s'", full_path, source_id)

        # 2) Parse file to text — Unstructured 우선, fallback은 기존 파서
        content = _parse_file_advanced(full_path)

        if not content or not content.strip():
            return {
                "source_id": source_id,
                "document_id": doc_id,
                "chunks_created": 0,
                "success": False,
                "errors": [{"error": f"No text content extracted from: {file_path}"}],
            }

        logger.info("Extracted %d chars from '%s'", len(content), file_path)

        # 2.5) Auto-recommend chunking strategy based on file type
        ext = os.path.splitext(file_path)[1].lower()
        recommended = _recommend_strategy(ext, chunking_strategy)
        if recommended != chunking_strategy:
            logger.info("Auto-recommended '%s' strategy for %s file (was '%s')", recommended, ext, chunking_strategy)
            chunking_strategy = recommended

        # 3) Delegate to ingest_document for chunk → embed → store
        return ingest_document(
            source_id=source_id,
            content=content,
            document_id=doc_id,
            chunking_strategy=chunking_strategy,
            chunking_config=chunking_config,
            embedding_model=embedding_model,
        )

    except Exception as e:
        logger.error("File ingestion failed for '%s': %s", file_path, str(e))
        return {
            "source_id": source_id,
            "document_id": doc_id,
            "chunks_created": 0,
            "success": False,
            "errors": [{"error": str(e)}],
        }
